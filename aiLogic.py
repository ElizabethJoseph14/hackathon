from pptx import Presentation
import os

def generate_ppt(topic, num_slides, keypoints):
    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = topic

    # Add content slides
    for i in range(min(num_slides, len(keypoints))):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.shapes.placeholders[1]

        title.text = f"{topic} - Part {i+1}"
        content.text = keypoints[i]

    # Save file
    output_file = "data/generated.pptx"
    prs.save(output_file)
    return output_file
