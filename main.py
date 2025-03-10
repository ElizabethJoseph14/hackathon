from fastapi import FastAPI, HTTPException
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List
import os
from pptx import Presentation

app = FastAPI()

class PresentationRequest(BaseModel):
    name: str
    category: str
    topic: str
    slides: int
    keypoints: List[str]

@app.post("/generate-presentation/")
async def generate_presentation(request: PresentationRequest):
    save_path = "data/generated.pptx"

    # Ensure the 'data' folder exists
    os.makedirs("data", exist_ok=True)

    # Create a new PowerPoint presentation
    prs = Presentation()

    # Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Title Slide Layout
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = request.topic
    subtitle.text = f"By {request.name}"

    # Create slides with key points
    for i, keypoint in enumerate(request.keypoints[:request.slides], start=1):
        slide_layout = prs.slide_layouts[1]  # Title and Content Layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = f"Slide {i}"
        content.text = keypoint

    # Save the presentation
    prs.save(save_path)

    return FileResponse(save_path, filename="Smart_Presentation.pptx")
